"""
Erstellt eine PowerPoint-Präsentation der ARL-Simulationsergebnisse
für alle 9 PROVIDER-Szenarien.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path
from io import BytesIO

import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ---------------------------------------------------------------------------
# Pfade
# ---------------------------------------------------------------------------
BASE = Path(__file__).parent.parent
ANALYSIS = BASE / "analysis"
RESULTS  = Path(__file__).parent / "data" / "all_scenarios_results.json"
OUT_PPTX = BASE / "analysis" / "PROVIDER_Simulation_Ergebnisse.pptx"

with open(RESULTS) as f:
    meta = json.load(f)

# ---------------------------------------------------------------------------
# Design-Konstanten
# ---------------------------------------------------------------------------
# Farben
C_DARK    = RGBColor(0x1A, 0x23, 0x3A)   # Dunkelblau (Hintergrund Titelfolie)
C_ACCENT  = RGBColor(0x00, 0xB0, 0xD8)   # Cyan-Akzent
C_ACCENT2 = RGBColor(0x00, 0xD4, 0x7F)   # Grün
C_WARN    = RGBColor(0xFF, 0x57, 0x22)   # Orange-Rot
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xF5, 0xF7, 0xFA)   # Hellgrau Hintergrund
C_TEXT    = RGBColor(0x1A, 0x23, 0x3A)   # Dunkel für helle Folien
C_SUBTLE  = RGBColor(0x6B, 0x7B, 0x8D)   # Grau für Untertitel

SCENARIO_COLORS = {
    "S1": RGBColor(0x4C, 0xAF, 0x50),
    "S2": RGBColor(0x21, 0x96, 0xF3),
    "S3": RGBColor(0xFF, 0x57, 0x22),
    "S4": RGBColor(0xFF, 0x98, 0x00),
    "S5": RGBColor(0x00, 0xBC, 0xD4),
    "S6": RGBColor(0x9C, 0x27, 0xB0),
    "S7": RGBColor(0xF4, 0x43, 0x36),
    "S8": RGBColor(0x00, 0x96, 0x88),
    "S9": RGBColor(0x79, 0x55, 0x48),
}
SCENARIOS_META = [
    ("S1", "Soja / Tierfutter",        "Landwirtschaft"),
    ("S2", "Halbleiter",               "Industrie / Elektronik"),
    ("S3", "Pharma / Wirkstoffe",      "Gesundheit"),
    ("S4", "Düngemittel / AdBlue",     "Landwirtschaft / Mobilität"),
    ("S5", "Wasseraufbereitung",       "KRITIS / Daseinsvorsorge"),
    ("S6", "Rechenzentren",            "Digitale Infrastruktur"),
    ("S7", "Seltene Erden",            "Industrie / Energiewende"),
    ("S8", "Seefracht",                "Logistik / Handel"),
    ("S9", "Unterwasserkabel",         "Digitale Infrastruktur"),
]

W  = Inches(13.333)   # Folienbreite 16:9
H  = Inches(7.5)      # Folienhöhe

# ---------------------------------------------------------------------------
# Hilfsfunktionen
# ---------------------------------------------------------------------------

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # komplett leer
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill: RGBColor | None = None,
             line: RGBColor | None = None, line_width: int = 0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text: str, left, top, width, height,
             font_size: int = 18, bold: bool = False, color: RGBColor = C_TEXT,
             align=PP_ALIGN.LEFT, wrap: bool = True, italic: bool = False,
             font_name: str = "Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font_name
    return txBox

def add_image(slide, path: Path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(str(path), left, top, width, height)
    else:
        slide.shapes.add_picture(str(path), left, top, width)

def slide_number(slide, prs: Presentation, n: int, total: int,
                 light: bool = False):
    color = C_SUBTLE if light else RGBColor(0x88, 0x99, 0xAA)
    add_text(slide, f"{n} / {total}",
             W - Inches(1.2), H - Inches(0.35), Inches(1.1), Inches(0.3),
             font_size=9, color=color, align=PP_ALIGN.RIGHT)

def header_bar(slide, title: str, subtitle: str = "", dark: bool = False):
    bg_color = C_DARK if dark else C_LIGHT
    fg_color = C_WHITE if dark else C_TEXT
    sub_color = C_ACCENT if dark else C_SUBTLE

    add_rect(slide, 0, 0, W, Inches(1.05), fill=bg_color)
    add_text(slide, title,
             Inches(0.4), Inches(0.08), W - Inches(0.8), Inches(0.55),
             font_size=22, bold=True, color=fg_color)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.6), W - Inches(0.8), Inches(0.38),
                 font_size=12, color=sub_color, italic=True)
    # Akzentlinie
    add_rect(slide, 0, Inches(1.05), W, Pt(3), fill=C_ACCENT)

def footer_bar(slide, text: str = "PROVIDER — ARL-Simulation 2026", light: bool = True):
    color = C_SUBTLE if light else RGBColor(0x55, 0x66, 0x77)
    add_text(slide, text,
             Inches(0.4), H - Inches(0.35), Inches(8), Inches(0.3),
             font_size=8, color=color, italic=True)

# ---------------------------------------------------------------------------
# Präsentation aufbauen
# ---------------------------------------------------------------------------
prs = new_prs()
TOTAL = 15

print("Erstelle Präsentation ...")

# ============================================================
# Folie 1 — Titelfolie
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_DARK)

# Oberes Farbband
add_rect(slide, 0, 0, W, Inches(0.18), fill=C_ACCENT)

# Zentrales Weißfeld
add_rect(slide, Inches(0.6), Inches(1.3), Inches(12.1), Inches(4.0),
         fill=RGBColor(0x22, 0x2E, 0x4A))

# Akzentlinie links
add_rect(slide, Inches(0.6), Inches(1.3), Inches(0.07), Inches(4.0), fill=C_ACCENT)

# Titel
add_text(slide, "ARL-Simulation",
         Inches(1.0), Inches(1.6), Inches(11.5), Inches(1.0),
         font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "Alle 9 PROVIDER-Szenarien im Vergleich",
         Inches(1.0), Inches(2.5), Inches(11.5), Inches(0.7),
         font_size=22, bold=False, color=C_ACCENT, align=PP_ALIGN.LEFT)

# Trennlinie
add_rect(slide, Inches(1.0), Inches(3.25), Inches(10.5), Pt(1),
         fill=RGBColor(0x44, 0x55, 0x70))

# Untertitel
add_text(slide, "Dummy-Baseline  ·  365 Ticks  ·  Seed 42  ·  9 Szenarien",
         Inches(1.0), Inches(3.4), Inches(11.0), Inches(0.5),
         font_size=14, color=RGBColor(0xAA, 0xBB, 0xCC), italic=True)
add_text(slide, "Adversarial Resilience Learning — Versorgungssicherheits-Simulation",
         Inches(1.0), Inches(3.9), Inches(11.0), Inches(0.45),
         font_size=13, color=RGBColor(0x88, 0x99, 0xAA))

# Unten: Projekt + Datum
add_rect(slide, 0, Inches(6.6), W, Inches(0.9), fill=RGBColor(0x11, 0x1A, 0x2A))
add_text(slide, "PROVIDER — Proaktive Versorgungssicherheit durch dynamische Simulation",
         Inches(0.5), Inches(6.65), Inches(9), Inches(0.4),
         font_size=11, color=RGBColor(0x88, 0x99, 0xAA))
add_text(slide, "Februar 2026",
         Inches(11.0), Inches(6.65), Inches(2.0), Inches(0.4),
         font_size=11, color=C_ACCENT, align=PP_ALIGN.RIGHT)

# Szenario-Farbpunkte dekorativ
for i, (sid, _, _) in enumerate(SCENARIOS_META):
    x = Inches(1.0 + i * 1.25)
    add_rect(slide, x, Inches(5.8), Inches(0.9), Inches(0.22),
             fill=SCENARIO_COLORS[sid])
    add_text(slide, sid, x, Inches(6.0), Inches(0.9), Inches(0.3),
             font_size=9, color=C_WHITE, align=PP_ALIGN.CENTER, bold=True)

print("  Folie 1: Titel")

# ============================================================
# Folie 2 — Agenda
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Agenda", dark=False)
footer_bar(slide)
slide_number(slide, prs, 2, TOTAL)

items = [
    ("01", "Methodik & Setup",             "Simulationsarchitektur, Parameter, Dummy-Baseline"),
    ("02", "Szenarien-Übersicht",          "9 Lieferketten-Szenarien, Entities & Events"),
    ("03", "Health-Zeitverläufe",          "Systemstabilität über 365 Simulationstage"),
    ("04", "Resilienzkennzahlen",          "Finale Health, Worst-Case, Einbruchgeschwindigkeit"),
    ("05", "Event-Kaskaden",               "Persistenz und Ausbreitung von Störereignissen"),
    ("06", "Kritischste Entities",         "Schwächste Glieder in jeder Lieferkette"),
    ("07", "Radar-Analyse",                "5-dimensionaler Resilienzbenchmark"),
    ("08", "Implikationen für ARL",        "Trainingspriorität, Defender-Baseline, nächste Schritte"),
]

for i, (num, title, desc) in enumerate(items):
    y = Inches(1.25 + i * 0.72)
    add_rect(slide, Inches(0.4), y + Inches(0.05), Inches(0.45), Inches(0.45),
             fill=C_ACCENT)
    add_text(slide, num, Inches(0.4), y + Inches(0.05), Inches(0.45), Inches(0.45),
             font_size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, title, Inches(1.0), y, Inches(3.8), Inches(0.55),
             font_size=14, bold=True, color=C_TEXT)
    add_text(slide, desc, Inches(4.9), y + Inches(0.08), Inches(8.0), Inches(0.42),
             font_size=11, color=C_SUBTLE)
    if i < len(items) - 1:
        add_rect(slide, Inches(0.4), y + Inches(0.57), Inches(12.5), Pt(0.5),
                 fill=RGBColor(0xDD, 0xE3, 0xEA))

print("  Folie 2: Agenda")

# ============================================================
# Folie 3 — Methodik & Setup
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Methodik & Setup", "Simulationsarchitektur und Parametrisierung")
footer_bar(slide)
slide_number(slide, prs, 3, TOTAL)

# Linke Spalte: Architektur
add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.8), Inches(5.7),
         fill=C_WHITE, line=RGBColor(0xDD, 0xE3, 0xEA), line_width=1)
add_text(slide, "Simulationsarchitektur", Inches(0.6), Inches(1.3), Inches(5.4), Inches(0.4),
         font_size=13, bold=True, color=C_DARK)
arch_lines = [
    ("PDL-Szenario (YAML)", "Eingabe: Entities, Events, Kaskaden"),
    ("ProviderEnvironment", "palaestrAI-kompatibles Gym"),
    ("5-Phasen-Simulationsengine", "Actions → Events → Impact → Flow → Health"),
    ("DummyBrain / DummyMuscle", "Null-Steuerung = natürliche Baseline"),
    ("step_dict() API", "Standalone ohne Orchestrator"),
]
for i, (title, desc) in enumerate(arch_lines):
    y = Inches(1.8 + i * 0.9)
    add_rect(slide, Inches(0.55), y, Inches(0.06), Inches(0.35), fill=C_ACCENT)
    add_text(slide, title, Inches(0.75), y - Inches(0.02), Inches(5.2), Inches(0.35),
             font_size=11, bold=True, color=C_DARK)
    add_text(slide, desc, Inches(0.75), y + Inches(0.3), Inches(5.2), Inches(0.35),
             font_size=9, color=C_SUBTLE, italic=True)

# Rechte Spalte: Parameter
add_rect(slide, Inches(6.5), Inches(1.2), Inches(6.4), Inches(5.7),
         fill=C_WHITE, line=RGBColor(0xDD, 0xE3, 0xEA), line_width=1)
add_text(slide, "Experiment-Parameter", Inches(6.7), Inches(1.3), Inches(6.0), Inches(0.4),
         font_size=13, bold=True, color=C_DARK)

params = [
    ("Szenarien",         "9 (S1–S9, alle PROVIDER-Lieferketten)"),
    ("Ticks / Episode",   "365 (1 Simulationsjahr)"),
    ("Seed",              "42 (reproduzierbar)"),
    ("Agentensteuerung",  "Dummy (Null-Aktionen = Baseline)"),
    ("Reward-Struktur",   "Zero-Sum: Attacker + Defender = 1.0"),
    ("Entities gesamt",   f"Ø {np.mean([meta[s]['n_entities'] for s in meta]):.0f} pro Szenario"),
    ("Events gesamt",     f"Ø {np.mean([meta[s]['n_events'] for s in meta]):.0f} pro Szenario"),
    ("Health-Formel",     "0.5·supply + 0.3·(1/price) + 0.2·demand"),
]
for i, (key, val) in enumerate(params):
    y = Inches(1.8 + i * 0.6)
    add_text(slide, key, Inches(6.7), y, Inches(2.6), Inches(0.4),
             font_size=11, bold=True, color=C_TEXT)
    add_text(slide, val, Inches(9.4), y, Inches(3.3), Inches(0.4),
             font_size=11, color=C_SUBTLE)
    if i < len(params) - 1:
        add_rect(slide, Inches(6.6), y + Inches(0.42), Inches(6.1), Pt(0.5),
                 fill=RGBColor(0xEE, 0xF1, 0xF5))

print("  Folie 3: Methodik")

# ============================================================
# Folie 4 — Szenarien-Übersicht (Tabelle)
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Szenarien-Übersicht", "9 Lieferketten-Szenarien des PROVIDER-Projekts")
footer_bar(slide)
slide_number(slide, prs, 4, TOTAL)

# Tabellenkopf
headers = ["ID", "Szenario", "Sektor", "Entities", "Events", "H@365", "H_min"]
col_x   = [Inches(0.35), Inches(0.95), Inches(4.0), Inches(7.5), Inches(8.6),
           Inches(9.7), Inches(11.1)]
col_w   = [Inches(0.55), Inches(3.0), Inches(3.4), Inches(1.0), Inches(1.0),
           Inches(1.3), Inches(1.3)]

add_rect(slide, Inches(0.35), Inches(1.15), Inches(12.6), Inches(0.42),
         fill=C_DARK)
for hdr, x, w in zip(headers, col_x, col_w):
    add_text(slide, hdr, x + Inches(0.05), Inches(1.18), w, Inches(0.38),
             font_size=11, bold=True, color=C_WHITE)

for i, (sid, name, sector) in enumerate(SCENARIOS_META):
    r = meta[sid]
    y = Inches(1.57 + i * 0.55)
    row_fill = RGBColor(0xF0, 0xF4, 0xF8) if i % 2 == 0 else C_WHITE
    add_rect(slide, Inches(0.35), y, Inches(12.6), Inches(0.53), fill=row_fill)

    # Farbpunkt
    add_rect(slide, col_x[0] + Inches(0.05), y + Inches(0.12),
             Inches(0.28), Inches(0.28), fill=SCENARIO_COLORS[sid])
    add_text(slide, sid, col_x[0] + Inches(0.06), y + Inches(0.1),
             Inches(0.28), Inches(0.3),
             font_size=9, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    vals = [name, sector,
            str(r["n_entities"]), str(r["n_events"]),
            f"{r['mean_health_t365']:.3f}", f"{r['health_min']:.3f}"]
    for j, (val, x, w) in enumerate(zip(vals, col_x[1:], col_w[1:])):
        color = C_TEXT
        bold  = False
        if j == 4:  # H@365
            h = r["mean_health_t365"]
            color = RGBColor(0x00, 0x8A, 0x00) if h >= 0.72 else (
                    RGBColor(0xFF, 0x88, 0x00) if h >= 0.65 else C_WARN)
            bold = True
        add_text(slide, val, x + Inches(0.05), y + Inches(0.1), w, Inches(0.35),
                 font_size=10, color=color, bold=bold)

# Legende
add_text(slide, "Finale Health:  ", Inches(0.35), Inches(6.85), Inches(1.5), Inches(0.35),
         font_size=9, color=C_SUBTLE)
for lbl, col in [("≥ 0.72 stabil", RGBColor(0x00, 0x8A, 0x00)),
                  ("0.65–0.72 moderat", RGBColor(0xFF, 0x88, 0x00)),
                  ("< 0.65 kritisch", C_WARN)]:
    pass  # inline unten
add_text(slide, "■ ≥ 0.72 stabil   ■ 0.65–0.72 moderat   ■ < 0.65 kritisch",
         Inches(1.7), Inches(6.85), Inches(11), Inches(0.35),
         font_size=9, color=C_SUBTLE)

print("  Folie 4: Szenarien-Übersicht")

# ============================================================
# Folie 5 — Health-Zeitverläufe (9er Grid)
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Health-Zeitverläufe", "Mittlere System-Health je Szenario über 365 Ticks")
footer_bar(slide)
slide_number(slide, prs, 5, TOTAL)

img = ANALYSIS / "01_health_timeseries.png"
add_image(slide, img, Inches(0.3), Inches(1.1), Inches(12.7))

print("  Folie 5: Health-Zeitverläufe")

# ============================================================
# Folie 6 — Überlagerung aller Szenarien
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Alle Szenarien im Vergleich",
           "Überlagerte Zeitverläufe — zwei Cluster erkennbar")
footer_bar(slide)
slide_number(slide, prs, 6, TOTAL)

add_image(slide, ANALYSIS / "06_all_scenarios_overlay.png",
          Inches(0.3), Inches(1.1), Inches(8.7))

# Kommentarbox rechts
add_rect(slide, Inches(9.2), Inches(1.2), Inches(3.9), Inches(5.7),
         fill=C_WHITE, line=RGBColor(0xDD, 0xE3, 0xEA), line_width=1)
add_text(slide, "Zwei Cluster", Inches(9.4), Inches(1.35), Inches(3.5), Inches(0.4),
         font_size=13, bold=True, color=C_DARK)

cluster_a = [("S8", "Seefracht",       "0.782"),
             ("S5", "Wasser",          "0.716"),
             ("S7", "Seltene Erden",   "0.714"),
             ("S1", "Soja",            "0.701")]
cluster_b = [("S4", "Düngemittel",     "0.676"),
             ("S2", "Halbleiter",      "0.652"),
             ("S3", "Pharma",          "0.621"),
             ("S9", "Unterwasser",     "0.594"),
             ("S6", "Rechenzentren",   "0.569")]

add_rect(slide, Inches(9.3), Inches(1.8), Inches(3.7), Inches(0.28),
         fill=RGBColor(0xE8, 0xF5, 0xE9))
add_text(slide, "Klasse A — Moderat stabil (> 0.70)",
         Inches(9.35), Inches(1.82), Inches(3.6), Inches(0.25),
         font_size=9, bold=True, color=RGBColor(0x1B, 0x5E, 0x20))
for i, (sid, lbl, h) in enumerate(cluster_a):
    y = Inches(2.12 + i * 0.38)
    add_rect(slide, Inches(9.3), y + Inches(0.06), Inches(0.14), Inches(0.18),
             fill=SCENARIO_COLORS[sid])
    add_text(slide, f"{sid} {lbl}", Inches(9.5), y, Inches(2.5), Inches(0.35),
             font_size=10, color=C_TEXT)
    add_text(slide, h, Inches(12.1), y, Inches(0.9), Inches(0.35),
             font_size=10, bold=True, color=RGBColor(0x00, 0x8A, 0x00),
             align=PP_ALIGN.RIGHT)

add_rect(slide, Inches(9.3), Inches(3.72), Inches(3.7), Inches(0.28),
         fill=RGBColor(0xFB, 0xE9, 0xE7))
add_text(slide, "Klasse B — Kritisch (< 0.68)",
         Inches(9.35), Inches(3.74), Inches(3.6), Inches(0.25),
         font_size=9, bold=True, color=RGBColor(0xB7, 0x1C, 0x1C))
for i, (sid, lbl, h) in enumerate(cluster_b):
    y = Inches(4.04 + i * 0.38)
    add_rect(slide, Inches(9.3), y + Inches(0.06), Inches(0.14), Inches(0.18),
             fill=SCENARIO_COLORS[sid])
    add_text(slide, f"{sid} {lbl}", Inches(9.5), y, Inches(2.5), Inches(0.35),
             font_size=10, color=C_TEXT)
    add_text(slide, h, Inches(12.1), y, Inches(0.9), Inches(0.35),
             font_size=10, bold=True, color=C_WARN, align=PP_ALIGN.RIGHT)

print("  Folie 6: Überlagerung")

# ============================================================
# Folie 7 — Resilienzkennzahlen (3 Balkendiagramme)
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Resilienzkennzahlen im Vergleich",
           "Finale Health · Worst-Case · Einbruchgeschwindigkeit")
footer_bar(slide)
slide_number(slide, prs, 7, TOTAL)

add_image(slide, ANALYSIS / "02_resilience_comparison.png",
          Inches(0.3), Inches(1.1), Inches(12.7))

print("  Folie 7: Resilienzkennzahlen")

# ============================================================
# Folie 8 — Health-Heatmap
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "System-Health-Heatmap",
           "Alle Szenarien × Zeitachse — Rot=kritisch, Grün=stabil")
footer_bar(slide)
slide_number(slide, prs, 8, TOTAL)

add_image(slide, ANALYSIS / "03_health_heatmap.png",
          Inches(0.3), Inches(1.15), Inches(12.7))

# Annotierungen
insights = [
    (Inches(2.8), Inches(5.5), "S1: Grünphase bis T90"),
    (Inches(6.5), Inches(5.5), "S4: Zweistufiger\nEinbruch ~T184"),
    (Inches(10.2), Inches(5.5), "S8: Leichte\nErholung"),
]
for x, y, txt in insights:
    add_rect(slide, x - Inches(0.05), y - Inches(0.05),
             Inches(1.8), Inches(0.55), fill=RGBColor(0xFF, 0xFF, 0xCC),
             line=RGBColor(0xCC, 0xCC, 0x00), line_width=1)
    add_text(slide, txt, x, y, Inches(1.7), Inches(0.5),
             font_size=8, color=C_DARK, align=PP_ALIGN.CENTER)

print("  Folie 8: Health-Heatmap")

# ============================================================
# Folie 9 — Radar-Chart
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Radar-Analyse — 5 Resilienzdimensionen",
           "Normiert auf [0,1]; größere Fläche = resilienteres Szenario")
footer_bar(slide)
slide_number(slide, prs, 9, TOTAL)

add_image(slide, ANALYSIS / "07_radar_resilience.png",
          Inches(0.2), Inches(1.05), Inches(7.4))

# Erklärung der Dimensionen rechts
add_rect(slide, Inches(7.7), Inches(1.2), Inches(5.4), Inches(5.7),
         fill=C_WHITE, line=RGBColor(0xDD, 0xE3, 0xEA), line_width=1)
add_text(slide, "Dimensionen", Inches(7.9), Inches(1.35), Inches(5.0), Inches(0.38),
         font_size=13, bold=True, color=C_DARK)

dims = [
    ("Health @365",         "Finale mittlere System-Health"),
    ("Health Min",          "Schlechtester Tick (Worst Case)"),
    ("Einbruch (inv.)",     "Normierter Drop-Tick (spät = gut)"),
    ("Event-Persistenz (inv.)", "Anteil nicht-dauerhafter Events"),
    ("Defender Reward",     "Mittlerer Reward des Verteidigers"),
]
for i, (dim, desc) in enumerate(dims):
    y = Inches(1.85 + i * 0.88)
    add_rect(slide, Inches(7.8), y, Inches(0.06), Inches(0.7), fill=C_ACCENT)
    add_text(slide, dim,  Inches(8.0), y, Inches(4.8), Inches(0.38),
             font_size=11, bold=True, color=C_DARK)
    add_text(slide, desc, Inches(8.0), y + Inches(0.36), Inches(4.8), Inches(0.38),
             font_size=10, color=C_SUBTLE, italic=True)

add_rect(slide, Inches(7.8), Inches(6.25), Inches(5.1), Inches(0.5),
         fill=RGBColor(0xE3, 0xF2, 0xFD))
add_text(slide,
         "S8 Seefracht hat die größte Raderfläche → resilientestes Szenario",
         Inches(7.9), Inches(6.28), Inches(5.0), Inches(0.42),
         font_size=10, bold=True, color=RGBColor(0x0D, 0x47, 0xA1))

print("  Folie 9: Radar-Chart")

# ============================================================
# Folie 10 — Event-Kaskaden
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Event-Aktivierungsmuster",
           "Zeitliche Persistenz und Ausbreitung von Störereignissen")
footer_bar(slide)
slide_number(slide, prs, 10, TOTAL)

add_image(slide, ANALYSIS / "04_event_patterns.png",
          Inches(0.3), Inches(1.1), Inches(12.7))

print("  Folie 10: Event-Kaskaden")

# ============================================================
# Folie 11 — Kritischste Entities
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Kritischste Entities je Szenario",
           "Schwächste Glieder in jeder Lieferkette (finale Health)")
footer_bar(slide)
slide_number(slide, prs, 11, TOTAL)

add_image(slide, ANALYSIS / "08_worst_entities.png",
          Inches(0.3), Inches(1.1), Inches(8.7))

# Tabelle rechts: eine Zeile je Szenario
add_rect(slide, Inches(9.2), Inches(1.2), Inches(3.9), Inches(5.8),
         fill=C_WHITE, line=RGBColor(0xDD, 0xE3, 0xEA), line_width=1)
add_text(slide, "Kritischste Entity", Inches(9.4), Inches(1.3), Inches(3.5), Inches(0.38),
         font_size=12, bold=True, color=C_DARK)

worst_entities = [
    ("S1", "feed_mills",                   0.383),
    ("S2", "automotive_consumers",         0.433),
    ("S3", "india_alternative_sourcing",   0.388),
    ("S4", "slaughterhouses",              0.500),
    ("S5", "small_rural_utilities",        0.500),
    ("S6", "diesel_supply",               0.500),
    ("S7", "ndfeb_magnet_producers",       0.275),
    ("S8", "rotterdam_port",              0.500),
    ("S9", "financial_sector_de",         0.444),
]
for i, (sid, eid, h) in enumerate(worst_entities):
    y = Inches(1.75 + i * 0.53)
    add_rect(slide, Inches(9.25), y + Inches(0.08), Inches(0.12), Inches(0.28),
             fill=SCENARIO_COLORS[sid])
    name_short = eid.replace("_", " ")
    add_text(slide, f"{sid}  {name_short}", Inches(9.42), y, Inches(2.9), Inches(0.35),
             font_size=9, color=C_TEXT)
    hcolor = (RGBColor(0xB7, 0x1C, 0x1C) if h < 0.4
              else RGBColor(0xFF, 0x88, 0x00) if h < 0.5
              else C_SUBTLE)
    add_text(slide, f"{h:.3f}", Inches(12.2), y, Inches(0.75), Inches(0.35),
             font_size=10, bold=True, color=hcolor, align=PP_ALIGN.RIGHT)

add_rect(slide, Inches(9.25), Inches(6.62), Inches(3.7), Inches(0.35),
         fill=RGBColor(0xFD, 0xE9, 0xE7))
add_text(slide, "⚠ S7: NdFeB-Magnete kritischste Entity (0.275)",
         Inches(9.3), Inches(6.64), Inches(3.6), Inches(0.3),
         font_size=9, bold=True, color=RGBColor(0xB7, 0x1C, 0x1C))

print("  Folie 11: Kritischste Entities")

# ============================================================
# Folie 12 — Spotlight: Resilientestes vs. Kritischstes
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Spotlight: Extreme Fälle",
           "S8 Seefracht (resilientestes) vs. S6 & S9 (kritischste Szenarien)")
footer_bar(slide)
slide_number(slide, prs, 12, TOTAL)

# S8-Karte
add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.9), Inches(5.7),
         fill=RGBColor(0xE0, 0xF2, 0xF1), line=SCENARIO_COLORS["S8"], line_width=2)
add_rect(slide, Inches(0.4), Inches(1.2), Inches(5.9), Inches(0.45),
         fill=SCENARIO_COLORS["S8"])
add_text(slide, "S8 — Seefracht  ✓ Resilientestes Szenario",
         Inches(0.55), Inches(1.25), Inches(5.6), Inches(0.38),
         font_size=12, bold=True, color=C_WHITE)
s8_facts = [
    ("Finale Health",        "0.782  (Rang 1 von 9)"),
    ("Minimum Health",       "0.754  (höchste Untergrenze)"),
    ("Drop-Tick",            "T5  (moderater Einbruch)"),
    ("Event-Persistenz",     "82% dauerhaft aktiv"),
    ("Defender Reward Ø",    "0.775"),
    ("Charakteristik",       "Moderater Einbruch, leichte\nnatürliche Erholung ab T270"),
    ("ARL-Eignung",          "Ideal für erste Trainingsläufe"),
]
for i, (k, v) in enumerate(s8_facts):
    y = Inches(1.75 + i * 0.7)
    add_text(slide, k + ":", Inches(0.6), y, Inches(2.2), Inches(0.38),
             font_size=10, bold=True, color=RGBColor(0x00, 0x60, 0x5A))
    add_text(slide, v, Inches(2.85), y, Inches(3.3), Inches(0.5),
             font_size=10, color=C_TEXT)

# S6 + S9 Karte
add_rect(slide, Inches(6.55), Inches(1.2), Inches(6.4), Inches(5.7),
         fill=RGBColor(0xFB, 0xE9, 0xE7), line=C_WARN, line_width=2)
add_rect(slide, Inches(6.55), Inches(1.2), Inches(6.4), Inches(0.45), fill=C_WARN)
add_text(slide, "S6 & S9 — Rechenzentren / Unterwasserkabel  ⚠",
         Inches(6.7), Inches(1.25), Inches(6.1), Inches(0.38),
         font_size=12, bold=True, color=C_WHITE)

for row_i, (sid, facts) in enumerate([
    ("S6", [
        ("Finale Health",  "0.569  (Rang 9 von 9)"),
        ("Drop-Tick",      "T28  (langsamster Einbruch,\naber dauerhaft tief)"),
        ("Events",         "0% dauerhaft → zyklisches\nVerhalten"),
        ("Charakteristik", "Langsamer Start, dann\npermanent kritisch"),
    ]),
    ("S9", [
        ("Finale Health",  "0.594  (Rang 8 von 9)"),
        ("Drop-Tick",      "T2  (schnellster Kollaps\naller Szenarien)"),
        ("Kritisch",       "financial_sector_de: 0.444"),
        ("Charakteristik", "Sofortiger tiefer Einbruch,\nkein Recovery"),
    ])
]):
    x_off = Inches(6.7 + row_i * 3.1)
    add_rect(slide, x_off - Inches(0.1), Inches(1.75), Inches(2.9), Inches(0.3),
             fill=SCENARIO_COLORS[sid])
    add_text(slide, sid, x_off - Inches(0.05), Inches(1.78), Inches(2.8), Inches(0.25),
             font_size=10, bold=True, color=C_WHITE)
    for i, (k, v) in enumerate(facts):
        y = Inches(2.15 + i * 1.1)
        add_text(slide, k + ":", x_off - Inches(0.05), y, Inches(2.8), Inches(0.35),
                 font_size=9, bold=True, color=RGBColor(0x8B, 0x1A, 0x00))
        add_text(slide, v, x_off - Inches(0.05), y + Inches(0.35), Inches(2.8), Inches(0.55),
                 font_size=9, color=C_TEXT)

print("  Folie 12: Spotlight")

# ============================================================
# Folie 13 — Entity-Health-Heatmaps
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Entity-Health-Heatmaps",
           "Kaskadenstruktur innerhalb jeder Lieferkette")
footer_bar(slide)
slide_number(slide, prs, 13, TOTAL)

add_image(slide, ANALYSIS / "05_entity_health_heatmaps.png",
          Inches(0.3), Inches(1.1), Inches(12.7))

print("  Folie 13: Entity-Heatmaps")

# ============================================================
# Folie 14 — Implikationen für ARL-Training
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_LIGHT)
header_bar(slide, "Implikationen für das ARL-Training",
           "Prioritäten, Defender-Baseline und empfohlene nächste Schritte")
footer_bar(slide)
slide_number(slide, prs, 14, TOTAL)

# 3 Spalten
col_configs = [
    ("Trainings-\npriorität", C_ACCENT, [
        "S8 Seefracht als Einstiegsszenario (stabilste Baseline)",
        "S1 Soja: gradueller Verfall ideal für Curriculum Learning",
        "S4 Düngemittel: zweistufiger Einbruch als Fortgeschrittenen-Szenario",
        "S6/S9 erst nach Grundtraining (zu instabil für Exploration)",
    ]),
    ("Defender-\nBaseline", RGBColor(0x00, 0x96, 0x88), [
        "Ø Defender Reward ohne Training: 0.697",
        "Beste natürliche Baseline: S8 (0.775)",
        "Schwächste Baseline: S9 (0.601) → hoher Verbesserungsspielraum",
        "Zero-Sum garantiert messbaren Fortschritt",
    ]),
    ("Nächste\nSchritte", RGBColor(0x7B, 0x1F, 0xA2), [
        "PPO-Training für S8 mit train_ppo.py starten",
        "Vergleich: PPO-Defender vs. Dummy-Baseline",
        "Heatmap-Analyse nach Training (verbesserte Entities?)",
        "Multi-Szenario-Training für robuste Generalisierung",
    ]),
]

for ci, (title, color, points) in enumerate(col_configs):
    x = Inches(0.35 + ci * 4.35)
    add_rect(slide, x, Inches(1.2), Inches(4.1), Inches(0.5), fill=color)
    add_text(slide, title, x + Inches(0.1), Inches(1.23), Inches(3.9), Inches(0.5),
             font_size=13, bold=True, color=C_WHITE)
    add_rect(slide, x, Inches(1.7), Inches(4.1), Inches(5.15),
             fill=C_WHITE, line=RGBColor(0xDD, 0xE3, 0xEA), line_width=1)
    for i, point in enumerate(points):
        y = Inches(1.85 + i * 1.2)
        add_rect(slide, x + Inches(0.15), y + Inches(0.08),
                 Inches(0.08), Inches(0.28), fill=color)
        add_text(slide, point, x + Inches(0.35), y, Inches(3.6), Inches(1.1),
                 font_size=10, color=C_TEXT, wrap=True)

print("  Folie 14: Implikationen")

# ============================================================
# Folie 15 — Fazit & Ausblick
# ============================================================
slide = blank_slide(prs)
fill_bg(slide, C_DARK)

add_rect(slide, 0, 0, W, Inches(0.18), fill=C_ACCENT)
add_rect(slide, 0, Inches(7.32), W, Inches(0.18), fill=C_ACCENT)

add_text(slide, "Fazit & Ausblick",
         Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
         font_size=26, bold=True, color=C_WHITE)
add_rect(slide, Inches(0.6), Inches(1.05), Inches(11.8), Pt(1.5),
         fill=RGBColor(0x44, 0x55, 0x70))

# Kernaussagen
kernaussagen = [
    ("Baseline etabliert",
     "Alle 9 Szenarien vollständig simuliert — reproduzierbare Dummy-Baseline für ARL-Benchmarking"),
    ("Zwei Resilienzkategorien",
     "Obere Gruppe (S8/S5/S7/S1): H ≥ 0.70 · Untere Gruppe (S4/S2/S3/S9/S6): H < 0.68"),
    ("Kritischste Schwachstellen",
     "NdFeB-Magnete (S7, H=0.275), Futtermühlen (S1, H=0.383), India-Sourcing (S3, H=0.388)"),
    ("Event-Persistenz dominiert",
     "Ø 73% aller Events dauerhaft aktiv — Systeme finden kein natürliches Gleichgewicht"),
    ("Nächster Meilenstein",
     "PPO-Training gegen diese Baseline · Ziel: Defender schlägt Dummy signifikant"),
]

for i, (bold_part, rest) in enumerate(kernaussagen):
    y = Inches(1.2 + i * 1.1)
    add_rect(slide, Inches(0.5), y + Inches(0.1), Inches(0.12), Inches(0.5),
             fill=C_ACCENT)
    add_text(slide, bold_part, Inches(0.75), y, Inches(3.8), Inches(0.42),
             font_size=12, bold=True, color=C_ACCENT)
    add_text(slide, rest, Inches(0.75), y + Inches(0.4), Inches(11.8), Inches(0.55),
             font_size=11, color=RGBColor(0xCC, 0xD5, 0xE0), wrap=True)

# Footer
add_rect(slide, 0, Inches(6.6), W, Inches(0.72), fill=RGBColor(0x11, 0x1A, 0x2A))
add_text(slide, "PROVIDER — Proaktive Versorgungssicherheit durch dynamische Simulation  ·  BMFTR-gefördertes Verbundprojekt",
         Inches(0.5), Inches(6.7), Inches(11.5), Inches(0.35),
         font_size=9, color=RGBColor(0x88, 0x99, 0xAA), italic=True)
add_text(slide, f"Seite {TOTAL} / {TOTAL}",
         Inches(11.5), Inches(6.7), Inches(1.5), Inches(0.35),
         font_size=9, color=C_ACCENT, align=PP_ALIGN.RIGHT)

print("  Folie 15: Fazit")

# ---------------------------------------------------------------------------
# Speichern
# ---------------------------------------------------------------------------
prs.save(str(OUT_PPTX))
size_mb = OUT_PPTX.stat().st_size / 1_000_000
print(f"\n✓ Präsentation gespeichert: {OUT_PPTX}")
print(f"  {TOTAL} Folien · {size_mb:.1f} MB")
